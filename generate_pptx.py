#!/usr/bin/env python3
"""Generate presentation.pptx from the AI Agent Platform presentation content."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy

# ── Colours ──────────────────────────────────────────────────────────────────
BG      = RGBColor(0x08, 0x0c, 0x10)
SURFACE = RGBColor(0x0e, 0x14, 0x19)
SURF2   = RGBColor(0x14, 0x1c, 0x24)
BORDER  = RGBColor(0x1e, 0x28, 0x32)
TEXT    = RGBColor(0xdc, 0xe8, 0xf0)
MUTED   = RGBColor(0x60, 0x70, 0x80)
DIM     = RGBColor(0x3a, 0x4a, 0x5a)
BLUE    = RGBColor(0x3d, 0x9c, 0xf0)
GREEN   = RGBColor(0x34, 0xc9, 0x74)
PURPLE  = RGBColor(0xa4, 0x7c, 0xff)
ORANGE  = RGBColor(0xff, 0x96, 0x40)
RED     = RGBColor(0xff, 0x57, 0x57)
WHITE   = RGBColor(0xff, 0xff, 0xff)
CODE_BG = RGBColor(0x05, 0x08, 0x10)
RED_DIM = RGBColor(0x1a, 0x05, 0x05)
GREEN_DIM = RGBColor(0x0e, 0x2e, 0x1c)

# ── Slide dimensions (widescreen 16:9) ────────────────────────────────────────
W = Cm(33.867)
H = Cm(19.05)

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # completely blank
    slide  = prs.slides.add_slide(layout)
    return slide


def fill_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color=SURFACE, line_color=BORDER, line_width=Pt(0.5)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, font_size=Pt(12), bold=False, color=TEXT,
                font_name="Calibri", align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_label(slide, x, y, w, h, text, size=Pt(9), color=DIM, upper=True, bold=True):
    t = text.upper() if upper else text
    return add_textbox(slide, x, y, w, h, t, font_size=size, color=color, bold=bold)


def add_heading(slide, x, y, w, h, text, size=Pt(28), color=WHITE):
    return add_textbox(slide, x, y, w, h, text, font_size=size, bold=True,
                       color=color, font_name="Calibri Bold")


def add_chip(slide, x, y, text, bg_color=SURFACE, text_color=BLUE, border_color=BORDER):
    chip_w = Cm(3.2)
    chip_h = Cm(0.55)
    r = add_rect(slide, x, y, chip_w, chip_h, fill_color=bg_color, line_color=border_color)
    add_textbox(slide, x, y, chip_w, chip_h, text.upper(),
                font_size=Pt(7), bold=True, color=text_color,
                font_name="Courier New", align=PP_ALIGN.CENTER)
    return chip_w, chip_h


def add_code_block(slide, x, y, w, h, code_text):
    add_rect(slide, x, y, w, h, fill_color=CODE_BG, line_color=BORDER)
    tb = slide.shapes.add_textbox(x + Cm(0.3), y + Cm(0.2), w - Cm(0.6), h - Cm(0.4))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = code_text
    run.font.size = Pt(8)
    run.font.name = "Courier New"
    run.font.color.rgb = TEXT


def add_bullet_list(slide, x, y, w, items, color=MUTED, size=Pt(10), indent=Cm(0.4)):
    cur_y = y
    for item in items:
        # bullet dot
        dot = slide.shapes.add_textbox(x, cur_y, Cm(0.3), Cm(0.4))
        dp = dot.text_frame.paragraphs[0]
        dr = dp.add_run()
        dr.text = "›"
        dr.font.color.rgb = BLUE
        dr.font.size = size
        # text
        tb = slide.shapes.add_textbox(x + indent, cur_y, w - indent, Cm(0.45))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = item
        run.font.size = size
        run.font.color.rgb = color
        cur_y += Cm(0.5)
    return cur_y


def add_flow_box(slide, x, y, w, h, label, sub=None, border_color=BORDER):
    add_rect(slide, x, y, w, h, fill_color=SURFACE, line_color=border_color, line_width=Pt(1))
    if sub:
        add_textbox(slide, x, y + Cm(0.1), w, Cm(0.35), label,
                    font_size=Pt(9), bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Cm(0.45), w, Cm(0.5), sub,
                    font_size=Pt(8), color=MUTED, align=PP_ALIGN.CENTER)
    else:
        add_textbox(slide, x, y, w, h, label,
                    font_size=Pt(9), bold=True, color=TEXT, align=PP_ALIGN.CENTER)


def add_arrow(slide, x, y, color=BLUE):
    add_textbox(slide, x, y, Cm(0.5), Cm(0.4), "→",
                font_size=Pt(12), bold=True, color=color, align=PP_ALIGN.CENTER)


# ── Slides ────────────────────────────────────────────────────────────────────

def make_s0_cover(prs):
    slide = blank_slide(prs)
    fill_bg(slide)

    # Title
    add_textbox(slide, Cm(5), Cm(3.5), Cm(24), Cm(1.2),
                "VICI TAKE-HOME · MAR 2026",
                font_size=Pt(9), bold=True, color=BLUE, font_name="Courier New")

    add_heading(slide, Cm(5), Cm(5), Cm(24), Cm(3.5),
                "AI Agent\nPlatform", size=Pt(48))

    add_textbox(slide, Cm(5), Cm(9.2), Cm(22), Cm(1.2),
                "Three production-grade services built with strict TDD, deployed on Zeabur,\nand driven by Claude Code as the AI coding agent.",
                font_size=Pt(12), color=MUTED)

    # Chips row
    chips = [("Task 1 · LiteLLM Proxy", BLUE), ("Task 2 · GitHub Skills", PURPLE), ("Task 3 · Browser Agent", GREEN)]
    cx = Cm(5)
    for label, col in chips:
        add_rect(slide, cx, Cm(11.2), Cm(4.5), Cm(0.55),
                 fill_color=RGBColor(0x0e, 0x14, 0x19), line_color=col)
        add_textbox(slide, cx, Cm(11.2), Cm(4.5), Cm(0.55), label.upper(),
                    font_size=Pt(7.5), bold=True, color=col, font_name="Courier New", align=PP_ALIGN.CENTER)
        cx += Cm(4.8)

    # Divider
    line = slide.shapes.add_shape(1, Cm(5), Cm(12.4), Cm(24), Cm(0.02))
    line.fill.solid(); line.fill.fore_color.rgb = BORDER
    line.line.fill.background()

    # Meta
    meta = [("Tasks", "3 / 3"), ("Unit Tests", "60+"), ("Workflow", "TDD"), ("Deploy", "Zeabur"), ("AI Tool", "Claude Code")]
    mx = Cm(5)
    for label, val in meta:
        add_textbox(slide, mx, Cm(12.6), Cm(4), Cm(0.35), label.upper(),
                    font_size=Pt(7.5), color=DIM, font_name="Courier New")
        add_textbox(slide, mx, Cm(13.1), Cm(4), Cm(0.6), val,
                    font_size=Pt(14), bold=True, color=TEXT, font_name="Calibri Bold")
        mx += Cm(4.5)


def make_s1_architecture(prs):
    slide = blank_slide(prs)
    fill_bg(slide)

    add_textbox(slide, Cm(2), Cm(1.0), Cm(10), Cm(0.5), "OVERVIEW",
                font_size=Pt(8), bold=True, color=BLUE, font_name="Courier New")
    add_heading(slide, Cm(2), Cm(1.6), Cm(28), Cm(1.2), "Architecture", size=Pt(30))
    add_textbox(slide, Cm(2), Cm(2.9), Cm(26), Cm(0.6),
                "Three independent services wired together — each deployed on Zeabur, communicating via HTTP.",
                font_size=Pt(10.5), color=MUTED)

    # Flow diagram
    fy = Cm(4.1)
    fh = Cm(1.2)
    boxes = [
        (Cm(2),    Cm(3.8), "Task 2\nGitHub Skills", "/gh-test · /gh-lint\n/gh-deploy · /gh-status", PURPLE),
        (Cm(7.4),  Cm(3.8), "GitHub Actions", "test.yml · lint.yml\ndeploy.yml", ORANGE),
        (Cm(13.5), Cm(3.8), "Task 3\nBrowser Agent", "POST /run-task", GREEN),
        (Cm(18.9), Cm(3.8), "Task 1\nLiteLLM Proxy", "POST /v1/chat/completions", BLUE),
        (Cm(24.3), Cm(3.8), "Anthropic\nClaude CLI", 'claude -p "..."', RED),
    ]
    for bx, by, label, sub, col in boxes:
        add_rect(slide, bx, by, Cm(4.8), fh, fill_color=SURFACE, line_color=col, line_width=Pt(1.2))
        lines = label.split("\n")
        add_textbox(slide, bx, by + Cm(0.05), Cm(4.8), Cm(0.45), lines[0],
                    font_size=Pt(8.5), bold=True, color=col, align=PP_ALIGN.CENTER)
        add_textbox(slide, bx, by + Cm(0.5), Cm(4.8), Cm(0.6), sub,
                    font_size=Pt(7.5), color=MUTED, align=PP_ALIGN.CENTER)

    # Arrows
    for ax in [Cm(6.85), Cm(12.3), Cm(17.75), Cm(23.15)]:
        add_textbox(slide, ax, Cm(4.15), Cm(0.7), Cm(0.5), "→",
                    font_size=Pt(14), bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    # 3-column cards
    card_y = Cm(6.2)
    card_h = Cm(11.3)
    cards = [
        (Cm(2),    "GitHub Repo",
         ["github.com/freedom870601/claude-agent-platform", "Commit history shows strict TDD red→green pairs throughout"],
         BLUE),
        (Cm(12.5), "Live Services (Zeabur)",
         ["Task 1  litellm-proxy-vici.zeabur.app", "Task 3  browser-agent.zeabur.app", "Task 2  Claude Code only (no HTTP)"],
         GREEN),
        (Cm(23),   "AI Workflow",
         ["Entire project built with Claude Code", "Skills used: /project:tdd, /project:commit", "TDD loop driven by Claude — tests first, then impl", "All commits authored via Claude Code agent"],
         PURPLE),
    ]
    for cx, title, items, col in cards:
        cw = Cm(10)
        add_rect(slide, cx, card_y, cw, Cm(12.3), fill_color=SURFACE, line_color=BORDER)
        add_textbox(slide, cx + Cm(0.3), card_y + Cm(0.2), cw - Cm(0.6), Cm(0.4),
                    title.upper(), font_size=Pt(7.5), bold=True, color=DIM, font_name="Courier New")
        cy = card_y + Cm(0.8)
        for item in items:
            add_textbox(slide, cx + Cm(0.3), cy, cw - Cm(0.6), Cm(0.5), item,
                        font_size=Pt(9), color=col if items.index(item) == 0 else MUTED,
                        font_name="Courier New" if "zeabur" in item or "github.com" in item else "Calibri")
            cy += Cm(0.55)


def make_s2_task1(prs):
    slide = blank_slide(prs)
    fill_bg(slide)

    add_textbox(slide, Cm(2), Cm(1.0), Cm(10), Cm(0.5), "TASK 1",
                font_size=Pt(8), bold=True, color=BLUE, font_name="Courier New")
    add_heading(slide, Cm(2), Cm(1.6), Cm(28), Cm(1.4), "LiteLLM Proxy — Claude CLI as an API", size=Pt(26))

    # Left column
    lw = Cm(14)
    add_textbox(slide, Cm(2), Cm(3.4), lw, Cm(1),
                "Bridges the claude CLI into an OpenAI-compatible endpoint. Any client that speaks OpenAI can use Claude — no SDK changes.",
                font_size=Pt(10), color=MUTED)

    # Mini flow
    flow_y = Cm(4.7)
    for i, (label, col) in enumerate([("OpenAI request", BORDER), ("LiteLLM CustomLLM", BLUE), ("claude -p", RED)]):
        fx = Cm(2) + i * Cm(4.7)
        add_rect(slide, fx, flow_y, Cm(4.3), Cm(0.55), fill_color=SURFACE, line_color=col)
        add_textbox(slide, fx, flow_y, Cm(4.3), Cm(0.55), label,
                    font_size=Pt(8.5), color=TEXT, align=PP_ALIGN.CENTER, font_name="Courier New")
        if i < 2:
            add_textbox(slide, fx + Cm(4.3), flow_y + Cm(0.05), Cm(0.4), Cm(0.45), "→",
                        font_size=Pt(12), bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    add_textbox(slide, Cm(2), Cm(5.6), lw, Cm(0.35), "KEY IMPLEMENTATION DETAILS",
                font_size=Pt(7.5), bold=True, color=DIM, font_name="Courier New")

    bullets = [
        "ClaudeCLIProvider(CustomLLM) in custom_handler.py",
        "acompletion() + astreaming() async methods",
        "Multi-turn history → Human: …\\nAssistant: …",
        "Tools serialised into --system-prompt, JSON parsed back",
        "_SUBPROCESS_ENV = dict(os.environ) — avoids uvloop rejection",
        "120-second timeout for Claude CLI subprocess",
    ]
    add_bullet_list(slide, Cm(2), Cm(6.1), lw, bullets, size=Pt(9.5))

    # Right column — curl demo
    rw = Cm(15)
    rx = Cm(17)
    add_textbox(slide, rx, Cm(3.4), rw, Cm(0.35), "DEMO — CURL",
                font_size=Pt(7.5), bold=True, color=DIM, font_name="Courier New")

    code = (
        "# Non-streaming\n"
        "curl -X POST https://litellm-proxy-vici.zeabur.app/v1/chat/completions \\\n"
        "  -H \"Authorization: Bearer $KEY\" \\\n"
        "  -H \"Content-Type: application/json\" \\\n"
        "  -d '{\n"
        '    "model": "claude-cli/claude-sonnet-4-5",\n'
        '    "messages": [{"role":"user","content":"Hello"}]\n'
        "  }'\n\n"
        "# Streaming — returns SSE chunks\n"
        "curl ... -d '{..., \"stream\": true}'\n\n"
        "# Tool use\n"
        "curl ... -d '{\n"
        "  ...,\n"
        '  "tools": [{"type": "function", "function": {\n'
        '    "name": "get_weather",\n'
        '    "description": "Get current weather",\n'
        '    "parameters": { ... }\n'
        "  }}]\n"
        "}'"
    )
    add_code_block(slide, rx, Cm(3.9), rw, Cm(14.6), code)


def make_s3_task2(prs):
    slide = blank_slide(prs)
    fill_bg(slide)

    add_textbox(slide, Cm(2), Cm(1.0), Cm(10), Cm(0.5), "TASK 2",
                font_size=Pt(8), bold=True, color=PURPLE, font_name="Courier New")
    add_heading(slide, Cm(2), Cm(1.6), Cm(28), Cm(1.4),
                "GitHub Skills — CI/CD as Claude Commands", size=Pt(26))
    add_textbox(slide, Cm(2), Cm(3.3), Cm(28), Cm(0.6),
                "Four .md skill files in github-skills/skills/, symlinked into .claude/commands/. "
                "Invoke directly from Claude Code — no copy-paste, no context switching.",
                font_size=Pt(10), color=MUTED)

    # 4 skill cards
    skills = [
        ("/gh-status", "List recent workflow runs. Read-only — safe to call any time.", "repo (required) · limit (default: 10)", PURPLE, False),
        ("/gh-test",   "Trigger test.yml, poll until complete, return pass/fail + URL.", "repo (required) · branch (default: main)", PURPLE, False),
        ("/gh-lint",   "Trigger lint.yml, poll with 10s intervals (max 5 min).", "repo (required) · branch (default: main)", PURPLE, False),
        ("/gh-deploy", "Deploy with mandatory CONFIRM gate. Auto-detects repo & changed service.", "repo · branch · environment · service", RED, True),
    ]
    card_y = Cm(4.3)
    cw = Cm(7.5)
    for i, (cmd, desc, inp, col, danger) in enumerate(skills):
        cx = Cm(2) + i * Cm(7.9)
        add_rect(slide, cx, card_y, cw, Cm(4.0), fill_color=SURFACE,
                 line_color=col if danger else BORDER)
        add_textbox(slide, cx + Cm(0.3), card_y + Cm(0.2), cw, Cm(0.55), cmd,
                    font_size=Pt(13), bold=True, color=col, font_name="Courier New")
        add_textbox(slide, cx + Cm(0.3), card_y + Cm(0.9), cw - Cm(0.6), Cm(1.2), desc,
                    font_size=Pt(9), color=MUTED)
        add_textbox(slide, cx + Cm(0.3), card_y + Cm(2.5), cw - Cm(0.6), Cm(0.45), inp,
                    font_size=Pt(8), color=DIM, font_name="Courier New")

    # Bottom 2 columns
    col_y = Cm(9.0)
    # Safety design
    add_textbox(slide, Cm(2), col_y, Cm(15), Cm(0.35), "SAFETY DESIGN",
                font_size=Pt(7.5), bold=True, color=DIM, font_name="Courier New")
    safety_bullets = [
        "Deploy requires user to type CONFIRM before dispatch",
        "Validator warns on embedded credentials / Bearer tokens",
        "validator.py enforces schema: name, description, version, inputs, outputs",
        "test_all_skills_valid() runs in CI — all 4 files must pass",
    ]
    add_rect(slide, Cm(2), col_y + Cm(0.5), Cm(15), Cm(5), fill_color=SURFACE, line_color=BORDER)
    add_bullet_list(slide, Cm(2.4), col_y + Cm(0.8), Cm(14), safety_bullets, size=Pt(9.5))

    # Auto-detect code
    add_textbox(slide, Cm(18.5), col_y, Cm(14), Cm(0.35), "AUTO-DETECT IN /GH-DEPLOY",
                font_size=Pt(7.5), bold=True, color=DIM, font_name="Courier New")
    code = (
        "# Detect repo from git remote\n"
        "git remote get-url origin\n"
        "# → https://github.com/user/repo.git\n\n"
        "# Detect changed service from diff\n"
        "git diff --name-only HEAD~1\n"
        "# litellm-proxy/** → service=litellm-proxy\n"
        "# browser-agent/** → service=browser-agent\n"
        "# both / neither   → service=all\n\n"
        "# Symlink install\n"
        "ln -s ../github-skills/skills .claude/commands"
    )
    add_code_block(slide, Cm(18.5), col_y + Cm(0.5), Cm(14), Cm(5.2), code)


def make_s4_task3(prs):
    slide = blank_slide(prs)
    fill_bg(slide)

    add_textbox(slide, Cm(2), Cm(1.0), Cm(10), Cm(0.5), "TASK 3",
                font_size=Pt(8), bold=True, color=GREEN, font_name="Courier New")
    add_heading(slide, Cm(2), Cm(1.6), Cm(28), Cm(1.4),
                "Browser Agent — Natural Language → Web Actions", size=Pt(24))

    # Left: Tool-use loop
    lw = Cm(14)
    add_textbox(slide, Cm(2), Cm(3.3), lw, Cm(0.35), "TOOL-USE LOOP",
                font_size=Pt(7.5), bold=True, color=DIM, font_name="Courier New")

    steps = [
        ("U", BLUE,   "POST /run-task",                     '{"task": "What is the latest Python release?"}'),
        ("C", PURPLE, "Claude decides: call browser_search", "tool_calls → {name: \"browser_search\", args: {…}}"),
        ("T", GREEN,  "Execute: DuckDuckGo search",          "Result added back as tool message"),
        ("C", PURPLE, "Claude decides: navigate + extract_text", "Reads python.org/downloads body (≤8 000 chars)"),
        ("✓", GREEN,  "Claude calls done(answer)",           "Loop exits — returns answer + full logs"),
    ]
    loop_y = Cm(3.8)
    add_rect(slide, Cm(2), loop_y, lw, Cm(7.5), fill_color=SURFACE, line_color=BORDER)
    sy = loop_y + Cm(0.3)
    for num, col, title, detail in steps:
        circ = slide.shapes.add_shape(9, Cm(2.3), sy, Cm(0.55), Cm(0.55))  # oval
        circ.fill.solid(); circ.fill.fore_color.rgb = col
        circ.line.fill.background()
        add_textbox(slide, Cm(2.3), sy, Cm(0.55), Cm(0.55), num,
                    font_size=Pt(8), bold=True, color=BG, align=PP_ALIGN.CENTER)
        add_textbox(slide, Cm(3.2), sy, lw - Cm(1.5), Cm(0.35), title,
                    font_size=Pt(9.5), bold=True, color=TEXT)
        add_textbox(slide, Cm(3.2), sy + Cm(0.35), lw - Cm(1.5), Cm(0.3), detail,
                    font_size=Pt(8), color=DIM, font_name="Courier New")
        sy += Cm(1.35)

    # Right: tools + response
    rx = Cm(18)
    rw = Cm(14)
    add_textbox(slide, rx, Cm(3.3), rw, Cm(0.35), "6 BROWSER TOOLS",
                font_size=Pt(7.5), bold=True, color=DIM, font_name="Courier New")
    tool_bullets = [
        "browser_search(query) — DuckDuckGo HTML (no API key)",
        "navigate(url) — Playwright goto, returns title",
        "extract_text() — page body ≤ 8 000 chars",
        "click(selector) — Playwright click",
        "type_text(selector, text) — Playwright fill",
        "done(answer) — terminates loop with result",
    ]
    add_rect(slide, rx, Cm(3.8), rw, Cm(3.8), fill_color=SURFACE, line_color=BORDER)
    add_bullet_list(slide, rx + Cm(0.3), Cm(4.0), rw - Cm(0.5), tool_bullets, size=Pt(9.5))

    add_textbox(slide, rx, Cm(7.9), rw, Cm(0.35), "POST /RUN-TASK RESPONSE",
                font_size=Pt(7.5), bold=True, color=DIM, font_name="Courier New")
    code = (
        '{\n'
        '  "task_id": "uuid",\n'
        '  "status": "completed",\n'
        '  "result": "Python 3.13.2 released...",\n'
        '  "steps_taken": 3,\n'
        '  "logs": [{\n'
        '    "step": 1,\n'
        '    "action": "browser_search",\n'
        '    "input": {"query": "..."},\n'
        '    "output": "...",\n'
        '    "timestamp": "2026-03-15T..."\n'
        '  }]\n'
        '}'
    )
    add_code_block(slide, rx, Cm(8.4), rw, Cm(5.5), code)


def make_s5_tdd(prs):
    slide = blank_slide(prs)
    fill_bg(slide)

    add_textbox(slide, Cm(2), Cm(1.0), Cm(10), Cm(0.5), "QUALITY",
                font_size=Pt(8), bold=True, color=ORANGE, font_name="Courier New")
    add_heading(slide, Cm(2), Cm(1.6), Cm(28), Cm(1.2), "TDD Evidence — Red → Green", size=Pt(28))

    # Stats row
    stats = [("73+", "Unit Tests\n(+ 6 integration)", BLUE), ("28", "Task 1", BLUE), ("13", "Task 2", PURPLE), ("32", "Task 3", GREEN)]
    sx = Cm(2)
    sw = Cm(7.5)
    for num, label, col in stats:
        add_rect(slide, sx, Cm(3.4), sw, Cm(2.2), fill_color=SURFACE, line_color=BORDER)
        add_textbox(slide, sx, Cm(3.5), sw, Cm(1.1), num,
                    font_size=Pt(32), bold=True, color=col, font_name="Calibri Bold", align=PP_ALIGN.CENTER)
        add_textbox(slide, sx, Cm(4.6), sw, Cm(0.7), label,
                    font_size=Pt(8), color=MUTED, align=PP_ALIGN.CENTER)
        sx += Cm(7.8)

    add_textbox(slide, Cm(2), Cm(6.0), Cm(28), Cm(0.4), "ACTUAL GIT COMMIT PAIRS",
                font_size=Pt(7.5), bold=True, color=DIM, font_name="Courier New")

    pairs = [
        ("test(task1): add failing test for simple user message assembly",
         "feat(task1): implement assemble_prompt for simple user message"),
        ("test(task1): add failing tests for CLI command builder",
         "feat(task1): implement _build_cmd for claude CLI invocation"),
        ("test(task1): add failing test for astreaming yields text chunks",
         "feat(task1): implement astreaming with JSON line parsing"),
        ("test(task2): add failing tests for load_frontmatter, check_required_fields",
         "feat(task2): implement load_frontmatter, check_required_fields, validate_skill"),
        ("test(task2): add failing test that all skill files pass validate_skill",
         "feat(task2): add skill markdown files (gh-lint, gh-test, gh-deploy, gh-status)"),
        ("test(task3): add failing test for browser_search returns results list",
         "feat(task3): implement browser_search tool with DuckDuckGo HTML parser"),
    ]
    row_y = Cm(6.6)
    rh = Cm(1.7)
    for red_msg, green_msg in pairs:
        half = Cm(16)
        # Red cell
        add_rect(slide, Cm(2), row_y, half, rh, fill_color=RED_DIM, line_color=RED, line_width=Pt(0.5))
        add_textbox(slide, Cm(2.2), row_y + Cm(0.05), half - Cm(0.4), Cm(0.4),
                    "🔴 FAILING TEST", font_size=Pt(7), bold=True, color=RED, font_name="Courier New")
        add_textbox(slide, Cm(2.2), row_y + Cm(0.5), half - Cm(0.4), Cm(1.1),
                    red_msg, font_size=Pt(8), color=MUTED, font_name="Courier New")
        # Green cell
        add_rect(slide, Cm(2) + half, row_y, half, rh, fill_color=GREEN_DIM, line_color=GREEN, line_width=Pt(0.5))
        add_textbox(slide, Cm(2.2) + half, row_y + Cm(0.05), half - Cm(0.4), Cm(0.4),
                    "🟢 IMPLEMENTATION", font_size=Pt(7), bold=True, color=GREEN, font_name="Courier New")
        add_textbox(slide, Cm(2.2) + half, row_y + Cm(0.5), half - Cm(0.4), Cm(1.1),
                    green_msg, font_size=Pt(8), color=MUTED, font_name="Courier New")
        row_y += rh + Cm(0.08)


def make_s8_demo(prs):
    slide = blank_slide(prs)
    fill_bg(slide)

    add_textbox(slide, Cm(2), Cm(1.0), Cm(10), Cm(0.5), "LIVE",
                font_size=Pt(8), bold=True, color=GREEN, font_name="Courier New")
    add_heading(slide, Cm(2), Cm(1.6), Cm(28), Cm(1.2), "Interactive Demo", size=Pt(28))

    # Task 1 curl example (left)
    lw = Cm(15)
    add_textbox(slide, Cm(2), Cm(3.2), lw, Cm(0.4),
                "TASK 1 — LITELLM PROXY",
                font_size=Pt(8), bold=True, color=BLUE, font_name="Courier New")
    code1 = (
        "curl -X POST https://litellm-proxy-vici.zeabur.app\\\n"
        "         /v1/chat/completions \\\n"
        "  -H \"Authorization: Bearer $KEY\" \\\n"
        "  -H \"Content-Type: application/json\" \\\n"
        "  -d '{\n"
        '    "model": "claude-cli/claude-sonnet-4-5",\n'
        '    "messages": [{"role":"user","content":"Hello"}]\n'
        "  }'"
    )
    add_code_block(slide, Cm(2), Cm(3.7), lw, Cm(4.0), code1)

    # Task 3 curl example (right)
    rx = Cm(18.5)
    rw = Cm(13.5)
    add_textbox(slide, rx, Cm(3.2), rw, Cm(0.4),
                "TASK 3 — BROWSER AGENT",
                font_size=Pt(8), bold=True, color=GREEN, font_name="Courier New")
    code3 = (
        "curl -X POST https://browser-agent.zeabur.app/run-task \\\n"
        "  -H \"Authorization: Bearer $KEY\" \\\n"
        "  -H \"Content-Type: application/json\" \\\n"
        "  -d '{\n"
        '    "task": "What is the latest Python release?",\n'
        '    "max_steps": 10\n'
        "  }'"
    )
    add_code_block(slide, rx, Cm(3.7), rw, Cm(4.0), code3)

    # Task 2 skills badges
    add_textbox(slide, Cm(2), Cm(8.1), Cm(28), Cm(0.4),
                "TASK 2 — GITHUB SKILLS (run inside Claude Code only)",
                font_size=Pt(8), bold=True, color=PURPLE, font_name="Courier New")
    badges = ["/gh-status", "/gh-test", "/gh-lint", "/gh-deploy ⚠ requires CONFIRM"]
    bx = Cm(2)
    for badge in badges:
        col = RED if "deploy" in badge else PURPLE
        brd = RED if "deploy" in badge else RGBColor(0x3a, 0x2a, 0x60)
        bw = Cm(6.5) if "deploy" in badge else Cm(4.5)
        add_rect(slide, bx, Cm(8.7), bw, Cm(0.7),
                 fill_color=RGBColor(0x2e, 0x10, 0x10) if "deploy" in badge else RGBColor(0x23, 0x1a, 0x40),
                 line_color=brd)
        add_textbox(slide, bx, Cm(8.7), bw, Cm(0.7), badge,
                    font_size=Pt(11), bold=True, color=col, font_name="Courier New", align=PP_ALIGN.CENTER)
        bx += bw + Cm(0.4)

    # Big hyperlink button
    btn_y = Cm(10.2)
    btn_w = Cm(18)
    btn_x = Cm(8)
    add_rect(slide, btn_x, btn_y, btn_w, Cm(2.0),
             fill_color=RGBColor(0x0d, 0x2a, 0x40), line_color=BLUE, line_width=Pt(2))

    tb = slide.shapes.add_textbox(btn_x, btn_y, btn_w, Cm(2.0))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "▶  Open Live Demo in Browser  →"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = BLUE
    run.font.name = "Calibri Bold"
    # Add hyperlink
    rPr = run._r
    from pptx.oxml.ns import qn
    from lxml import etree
    hlinkClick = etree.SubElement(rPr, qn('a:hlinkClick'))
    # Add relationship
    rId = slide.part.relate_to(
        "presentation.html",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
        is_external=True
    )
    hlinkClick.set(qn('r:id'), rId)

    add_textbox(slide, Cm(2), Cm(12.5), Cm(28), Cm(0.5),
                "Note: The live demo requires running services and a LiteLLM master key. "
                "Open presentation.html in a browser for the interactive version.",
                font_size=Pt(9), color=DIM)


def make_s9_summary(prs):
    slide = blank_slide(prs)
    fill_bg(slide)

    add_textbox(slide, Cm(2), Cm(1.0), Cm(10), Cm(0.5), "SUMMARY",
                font_size=Pt(8), bold=True, color=GREEN, font_name="Courier New")
    add_heading(slide, Cm(2), Cm(1.6), Cm(28), Cm(1.2), "Tech Stack", size=Pt(30))

    # 2x2 cards
    stacks = [
        ("BACKEND",       BLUE,   ["Python 3.11", "FastAPI", "LiteLLM", "Playwright", "Pydantic v2", "uvicorn", "asyncio", "uv"]),
        ("AI / CLI",      ORANGE, ["Claude CLI", "claude-sonnet-4-5", "OpenAI tool_calls", "stream-json", "ANTHROPIC_API_KEY"]),
        ("CI/CD & DEPLOY",PURPLE, ["GitHub Actions", "gh CLI", "Docker", "Zeabur", "zeabur.yaml", "Claude Code Skills"]),
        ("TESTING",       GREEN,  ["pytest", "pytest-asyncio", "unittest.mock", "httpx TestClient", "73+ unit tests", "@pytest.mark.integration"]),
    ]
    positions = [(Cm(2), Cm(3.5)), (Cm(17.5), Cm(3.5)), (Cm(2), Cm(9.0)), (Cm(17.5), Cm(9.0))]
    cw = Cm(14.5)
    ch = Cm(5.0)
    for (cx, cy), (label, col, pills) in zip(positions, stacks):
        add_rect(slide, cx, cy, cw, ch, fill_color=SURFACE, line_color=BORDER)
        add_textbox(slide, cx + Cm(0.3), cy + Cm(0.15), cw, Cm(0.4), label,
                    font_size=Pt(8), bold=True, color=col, font_name="Courier New")
        # Pills in rows
        px = cx + Cm(0.3)
        py = cy + Cm(0.75)
        for pill in pills:
            pw = Cm(3.8)
            add_rect(slide, px, py, pw, Cm(0.45),
                     fill_color=SURF2, line_color=BORDER)
            add_textbox(slide, px, py, pw, Cm(0.45), pill,
                        font_size=Pt(8), color=MUTED, align=PP_ALIGN.CENTER)
            px += pw + Cm(0.25)
            if px + pw > cx + cw:
                px = cx + Cm(0.3)
                py += Cm(0.55)

    # Summary banner
    banner_y = Cm(14.6)
    add_rect(slide, Cm(2), banner_y, Cm(30), Cm(3.8),
             fill_color=GREEN_DIM, line_color=RGBColor(0x34, 0xc9, 0x74))
    add_textbox(slide, Cm(2.5), banner_y + Cm(0.3), Cm(20), Cm(0.8),
                "All three tasks complete ✓",
                font_size=Pt(18), bold=True, color=GREEN, font_name="Calibri Bold")
    add_textbox(slide, Cm(2.5), banner_y + Cm(1.2), Cm(20), Cm(0.5),
                "Strict TDD workflow · Zeabur deployment · GitHub CI/CD · Claude Skills symlinked",
                font_size=Pt(10), color=MUTED)

    for val, label, cx in [("73+", "Tests", Cm(23.5)), ("3", "Services", Cm(26.5)), ("4", "Skills", Cm(29.5))]:
        add_textbox(slide, cx, banner_y + Cm(0.4), Cm(2.5), Cm(1.0), val,
                    font_size=Pt(24), bold=True, color=GREEN, font_name="Calibri Bold", align=PP_ALIGN.CENTER)
        add_textbox(slide, cx, banner_y + Cm(1.5), Cm(2.5), Cm(0.4), label.upper(),
                    font_size=Pt(8), color=MUTED, font_name="Courier New", align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()
    print("Building slides...")
    make_s0_cover(prs)       ; print("  [1/8] Cover")
    make_s1_architecture(prs); print("  [2/8] Architecture")
    make_s2_task1(prs)       ; print("  [3/8] Task 1 — LiteLLM Proxy")
    make_s3_task2(prs)       ; print("  [4/8] Task 2 — GitHub Skills")
    make_s4_task3(prs)       ; print("  [5/8] Task 3 — Browser Agent")
    make_s5_tdd(prs)         ; print("  [6/8] TDD Evidence")
    make_s8_demo(prs)        ; print("  [7/8] Interactive Demo")
    make_s9_summary(prs)     ; print("  [8/8] Summary / Tech Stack")

    out = "presentation.pptx"
    prs.save(out)
    print(f"\n✓ Saved → {out}")


if __name__ == "__main__":
    main()
